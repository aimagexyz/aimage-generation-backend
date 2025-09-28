from aimage_supervision.settings import logger
import os
import tempfile

import httpx
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE as MSHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from aimage_supervision.clients.aws_s3 import get_s3_url_from_path

# Helper function to safely get values from annotation_data (dict or object)


def get_annotation_value(data, key, default=None):
    """Safely retrieves a value from annotation_data, whether it's a dict or object."""
    if isinstance(data, dict):
        return data.get(key, default)
    else:
        return getattr(data, key, default)


def get_page_number_from_subtask_name(subtask_name: str) -> int:
    # Extract page number safely
    try:
        page_number_str = subtask_name.split('slide')[1]
        if not page_number_str:
            logger.warning(
                f"Could not extract page number from subtask name: {subtask_name}")
            return 0
        return int(page_number_str)
    except Exception as e:
        logger.error(
            f"Error extracting page number from subtask name: {subtask_name}: {e}")
        return 0


def hex_to_rgb(hex_color):
    """Converts a hex color string (e.g., '#ff0000') to an RGBColor object."""
    hex_color = hex_color.lstrip('#')
    try:
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    except ValueError:
        logger.warning(
            f"Invalid hex color format: {hex_color}. Using default black.")
        return RGBColor(0, 0, 0)  # Default to black if conversion fails

# Helper function to check if two rectangles overlap


def rectangles_overlap(rect1, rect2):
    """
    检查两个矩形是否重叠
    rect1, rect2: {'left': int, 'top': int, 'width': int, 'height': int}
    """
    r1_right = rect1['left'] + rect1['width']
    r1_bottom = rect1['top'] + rect1['height']
    r2_right = rect2['left'] + rect2['width']
    r2_bottom = rect2['top'] + rect2['height']

    return not (r1_right <= rect2['left'] or
                r2_right <= rect1['left'] or
                r1_bottom <= rect2['top'] or
                r2_bottom <= rect1['top'])


def find_non_overlapping_position(slide, shape_left, shape_top, shape_width, shape_height):
    """
    为annotation找到一个不重叠的位置
    返回调整后的 (left, top) 坐标
    """
    proposed_rect = {
        'left': shape_left,
        'top': shape_top,
        'width': shape_width,
        'height': shape_height
    }

    # 获取现有的annotation形状
    existing_annotations = []
    for shape in slide.shapes:
        if shape.name.startswith('Annotation_'):
            existing_annotations.append({
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height
            })

    logger.info(
        f"Found {len(existing_annotations)} existing annotations when positioning new one")

    # 如果没有重叠，直接返回原位置
    overlap_found = any(rectangles_overlap(proposed_rect, existing)
                        for existing in existing_annotations)

    if not overlap_found:
        logger.info("No overlap detected, using original position")
        return shape_left, shape_top

    logger.info("Overlap detected, searching for alternative position")

    # 改进的偏移策略：使用多层偏移距离和更多方向
    base_offset = Inches(0.8)  # 基础偏移距离

    # 多层偏移尝试：从小到大的偏移距离
    for multiplier in [1, 1.5, 2, 2.5, 3]:
        offset_distance = Emu(base_offset * multiplier)

        # 更多的偏移方向，包括对角线和直线方向
        offset_attempts = [
            # 第一优先级：右下方向（通常不会遮挡内容）
            (offset_distance, offset_distance),       # 右下
            (offset_distance * 1.5, offset_distance),  # 更右下
            (offset_distance, offset_distance * 1.5),  # 更下右

            # 第二优先级：其他角落
            (offset_distance, -offset_distance),      # 右上
            (-offset_distance, offset_distance),      # 左下
            (-offset_distance, -offset_distance),     # 左上

            # 第三优先级：直线方向
            (offset_distance, 0),                     # 纯右
            (0, offset_distance),                     # 纯下
            (-offset_distance, 0),                    # 纯左
            (0, -offset_distance),                    # 纯上

            # 第四优先级：更大的对角线偏移
            (offset_distance * 2, offset_distance),   # 远右下
            (offset_distance, offset_distance * 2),   # 远下右
        ]

        for offset_x, offset_y in offset_attempts:
            new_left = Emu(shape_left + offset_x)
            new_top = Emu(shape_top + offset_y)

            # 确保新位置在幻灯片范围内（简单检查）
            if new_left < 0 or new_top < 0:
                continue

            new_rect = {
                'left': new_left,
                'top': new_top,
                'width': shape_width,
                'height': shape_height
            }

            # 检查新位置是否还有重叠
            new_overlap = any(rectangles_overlap(new_rect, existing)
                              for existing in existing_annotations)

            if not new_overlap:
                logger.info(
                    f"Found non-overlapping position with offset ({offset_x}, {offset_y}) at multiplier {multiplier}")
                return new_left, new_top

    # 如果所有偏移都还有重叠，使用系统化的网格布局
    logger.warning(
        "Could not find position with offset methods, using grid fallback")

    # 网格布局作为最后的fallback
    grid_size = Inches(1.2)  # 网格间距
    for row in range(5):  # 尝试5行
        for col in range(5):  # 尝试5列
            grid_left = Emu(shape_left + (col * grid_size))
            grid_top = Emu(shape_top + (row * grid_size))

            if grid_left < 0 or grid_top < 0:
                continue

            grid_rect = {
                'left': grid_left,
                'top': grid_top,
                'width': shape_width,
                'height': shape_height
            }

            grid_overlap = any(rectangles_overlap(grid_rect, existing)
                               for existing in existing_annotations)

            if not grid_overlap:
                logger.info(
                    f"Found non-overlapping position using grid layout at ({col}, {row})")
                return grid_left, grid_top

    # 实在找不到位置，就放在一个固定的远离位置
    logger.warning(
        "Could not find any non-overlapping position, using fallback position")
    return Emu(shape_left + Inches(3)), Emu(shape_top + Inches(2))


# Helper function to find the main asset image on the slide using a top-left priority rule.


def find_largest_image(slide, slide_width, slide_height):
    """
    Finds the main asset image on the slide using a top-left priority rule.
    This logic is based on the principle used in PPTX extraction.
    """
    main_asset_candidate = None
    min_left, min_top = float('inf'), float('inf')

    TOLERANCE_RATIO = 0.15  # 15%的误差容忍度
    margin = slide_width * TOLERANCE_RATIO if slide_width else 0

    # Heuristic to filter out small images (e.g., logos)
    slide_area = slide_width * slide_height if slide_width > 0 and slide_height > 0 else 0

    for shape in slide.shapes:
        is_picture = (
            shape.shape_type == MSO_SHAPE_TYPE.PICTURE or
            (shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, 'image'))
        )
        if not is_picture:
            continue

        if not (hasattr(shape, 'left') and hasattr(shape, 'top')):
            continue

        # Heuristic: filter out small images that are likely logos or decorative elements.
        if slide_area > 0 and (shape.width * shape.height / slide_area) < 0.03:
            logger.info(
                "Skipping small image on slide - likely a logo or decorative element."
            )
            continue

        # Top-left priority logic:
        # 1. If a new candidate is significantly to the left of the current best, it becomes the new best.
        # 2. If it's in the same vertical "column" (within the margin), the higher one wins.
        is_in_new_left_column = shape.left < (min_left - margin)
        is_in_same_column_and_higher = abs(
            shape.left - min_left) <= margin and shape.top < min_top

        if is_in_new_left_column or is_in_same_column_and_higher:
            min_left, min_top = shape.left, shape.top
            main_asset_candidate = shape

    if not main_asset_candidate:
        logger.warning("Could not find a main asset image on the slide.")

    return main_asset_candidate


# Updated function


async def add_comment_box_to_slide(slide, annotation_data, slide_width, slide_height):
    """
    Adds an annotation shape or a comment box to the slide based on annotation_data.

    Annotations ('annotation' type) are positioned relative to the largest image
    found on the slide, using 'rect' data. Comments ('comment' type) are placed
    at a fixed position (bottom).
    A visual marker (e.g., rectangle, oval) is drawn for 'annotation' types
    at the exact 'rect' location on the image. The associated text box is placed nearby.
    """
    if not isinstance(annotation_data, dict):
        logger.warning(f"Invalid annotation_data format: {annotation_data}")
        return

    # Use helper function to access data
    annotation_type = get_annotation_value(annotation_data, 'type')
    annotation_text = str(get_annotation_value(annotation_data, 'text', ''))
    annotation_color_hex = get_annotation_value(annotation_data, 'color')
    annotation_tool = get_annotation_value(annotation_data, 'tool')
    annotation_rect = get_annotation_value(annotation_data, 'rect')
    annotation_id = get_annotation_value(
        annotation_data, 'id', 'unknown')  # Get ID early

    # --- Variables for the text box / comment box ---
    # Default values, might be overridden based on type (comment vs annotation)
    current_shape_left = Inches(1.0)
    current_shape_top = Inches(1.0)
    current_shape_width = Inches(3.0)
    current_shape_height = Inches(1.5)
    current_shape_type = MSHAPE.ROUNDED_RECTANGLE
    current_fill_color = RGBColor(239, 157, 83)  # Default for comments
    current_line_color = RGBColor(150, 100, 50)  # Default for comments
    # Default for comments (white text)
    current_text_color = RGBColor(255, 255, 255)
    current_shape_name = f"Shape_{annotation_id}"

    is_comment = annotation_type == 'comment'

    if is_comment:
        # --- Handle Comments ---
        logger.debug(f"Adding comment: {annotation_text}")
        current_shape_height = Inches(1.5)
        current_shape_width = slide_width * 0.7
        current_shape_left = Inches(1.0)
        existing_comment_heights = 0
        for shape_in_slide in slide.shapes:
            if shape_in_slide.name.startswith("CommentBox_"):
                existing_comment_heights += shape_in_slide.height + Emu(50000)
        current_shape_top = slide_height - current_shape_height - \
            Inches(0.5) - existing_comment_heights
        current_shape_name = f"CommentBox_{annotation_id}"
        # Colors and shape type are already defaulted correctly for comments

    else:  # --- Handle Annotations ---
        logger.debug(f"Adding annotation: {annotation_text}")

        # Standard appearance for the annotation's TEXT BOX
        text_box_shape_type = MSHAPE.ROUNDED_RECTANGLE
        text_box_fill_color = RGBColor(239, 157, 83)  # Standard orange
        # Standard darker orange border
        text_box_line_color = RGBColor(150, 100, 50)
        text_box_text_color = RGBColor(255, 255, 255)  # White text

        # Color for the VISUAL MARKER's border
        visual_marker_line_rgb = RGBColor(255, 0, 0)  # Default Red if no color
        if annotation_color_hex:
            try:
                visual_marker_line_rgb = hex_to_rgb(annotation_color_hex)
            except ValueError:
                logger.warning(
                    f"Invalid hex color '{annotation_color_hex}' for visual marker. Defaulting to red.")

        largest_image = find_largest_image(
            slide, slide_width, slide_height)

        # Initial estimated dimensions for the text box (for collision detection)
        # These will be refined by SHAPE_TO_FIT_TEXT later for the actual shape.
        # Use Emu for find_non_overlapping_position
        text_box_collision_width = Inches(3.0).emu
        # Use Emu for find_non_overlapping_position
        text_box_collision_height = Inches(1.5).emu

        # Default positions for text_box if visual marker can't be placed
        text_box_final_left = Inches(1.0).emu
        text_box_final_top = Inches(1.0).emu

        if largest_image and annotation_rect and all(k in annotation_rect for k in ('x', 'y', 'width', 'height')):
            try:
                img_left_emu = largest_image.left
                img_top_emu = largest_image.top
                current_image_width_emu = largest_image.width
                current_image_height_emu = largest_image.height

                # --- Dynamic DPI Calculation (Aspect-Ratio Aware) ---
                # This is a heuristic approach because we don't know the frontend's viewport size
                # or the image's natural dimensions. It "guesses" the scaling logic based on
                # the image's aspect ratio in the PowerPoint slide.

                # Assumption:
                # - For landscape images (width > height), the frontend scaled the image to a fixed WIDTH.
                # - For portrait images (height >= width), the frontend scaled the image to a fixed HEIGHT.

                # Baseline values are derived from a known-good template.
                BASELINE_DPI = 220.0
                BASELINE_WIDTH_EMU = 5000000.0  # Adjusted for better landscape accuracy
                # For height baseline, assume the original baseline image was landscape (e.g., 16:9 or 4:3)
                # If it was 16:9, height would be width * (9/16) = width * 0.5625
                # If it was 4:3, height would be width * (3/4) = width * 0.75
                # Using a conservative 4:3 ratio for better compatibility
                BASELINE_HEIGHT_EMU = BASELINE_WIDTH_EMU * 0.68  # Assuming 4:3 aspect ratio

                assumed_dpi = BASELINE_DPI  # Default fallback

                if current_image_width_emu > 0 and current_image_height_emu > 0:
                    is_landscape = (current_image_width_emu /
                                    current_image_height_emu) > 1.0

                    if is_landscape:
                        # Use width-based scaling for landscape images
                        logger.info(
                            "Using WIDTH-based scaling for landscape image.")
                        dynamic_dpi = BASELINE_DPI * \
                            (BASELINE_WIDTH_EMU / current_image_width_emu)
                        assumed_dpi = dynamic_dpi
                    else:
                        # Use height-based scaling for portrait or square images
                        logger.info(
                            "Using HEIGHT-based scaling for portrait/square image.")
                        dynamic_dpi = BASELINE_DPI * \
                            (BASELINE_HEIGHT_EMU / current_image_height_emu)
                        assumed_dpi = dynamic_dpi
                else:
                    logger.warning(
                        "Image has zero width or height, using default DPI.")

                # 1. --- Calculate and Add VISUAL ANNOTATION SHAPE ---
                visual_rect_x_pixels = annotation_rect['x']
                visual_rect_y_pixels = annotation_rect['y']
                visual_rect_width_pixels = annotation_rect['width']
                visual_rect_height_pixels = annotation_rect['height']

                visual_rect_x_emu = Inches(
                    visual_rect_x_pixels / assumed_dpi).emu
                visual_rect_y_emu = Inches(
                    visual_rect_y_pixels / assumed_dpi).emu
                visual_rect_width_emu = Inches(
                    visual_rect_width_pixels / assumed_dpi).emu
                visual_rect_height_emu = Inches(
                    visual_rect_height_pixels / assumed_dpi).emu

                # --- Detailed logging for visual marker positioning ---
                logger.info(
                    f"[VISUAL MARKER DEBUG] Annotation ID: {annotation_id}")
                logger.info(
                    f"[VISUAL MARKER DEBUG]   Raw annotation_rect (pixels): {annotation_rect}")
                logger.info(
                    f"[VISUAL MARKER DEBUG]   Largest image L/T (emu): {img_left_emu}/{img_top_emu}, W/H (emu): {largest_image.width}/{largest_image.height}")
                logger.info(
                    f"[VISUAL MARKER DEBUG]   Assumed (Dynamic) DPI: {assumed_dpi}")
                logger.info(
                    f"[VISUAL MARKER DEBUG]   Rect X/Y (emu from pixels): {visual_rect_x_emu}/{visual_rect_y_emu}")
                logger.info(
                    f"[VISUAL MARKER DEBUG]   Rect W/H (emu from pixels): {visual_rect_width_emu}/{visual_rect_height_emu}")
                # --- End detailed logging ---

                if visual_rect_width_emu > 0 and visual_rect_height_emu > 0:
                    visual_shape_actual_left = Emu(
                        img_left_emu + visual_rect_x_emu)
                    visual_shape_actual_top = Emu(
                        img_top_emu + visual_rect_y_emu)

                    # --- More detailed logging for final position ---
                    logger.info(
                        f"[VISUAL MARKER DEBUG]   Calculated Visual Shape Actual L/T (emu): {visual_shape_actual_left}/{visual_shape_actual_top}")
                    # --- End more detailed logging ---

                    visual_shape_type_on_slide = MSHAPE.RECTANGLE  # Default
                    if annotation_tool == 'rect':
                        visual_shape_type_on_slide = MSHAPE.RECTANGLE
                    elif annotation_tool == 'ellipse':
                        visual_shape_type_on_slide = MSHAPE.OVAL
                    # Add more tool mappings if needed for visual shape

                    logger.info(
                        f"Attempting to add visual shape: type={visual_shape_type_on_slide}, L={visual_shape_actual_left}, T={visual_shape_actual_top}, W={visual_rect_width_emu}, H={visual_rect_height_emu}")
                    new_visual_shape = slide.shapes.add_shape(
                        visual_shape_type_on_slide,
                        visual_shape_actual_left,
                        visual_shape_actual_top,
                        Emu(visual_rect_width_emu),
                        Emu(visual_rect_height_emu)
                    )
                    new_visual_shape.name = f"VisualMarker_{annotation_id}"
                    new_visual_shape.fill.background()
                    new_visual_shape.line.color.rgb = visual_marker_line_rgb
                    new_visual_shape.line.width = Pt(2.0)
                    logger.info(
                        f"Added visual annotation marker for {annotation_id} at exact image coordinates.")

                    # Initial candidate position for the text box (e.g., at the visual marker's top-left)
                    # This will be adjusted by find_non_overlapping_position.
                    initial_text_box_left = visual_shape_actual_left + \
                        visual_rect_width_emu + Inches(0.1).emu
                    initial_text_box_top = visual_shape_actual_top
                else:
                    logger.warning(
                        f"Visual annotation rect for {annotation_id} has non-positive width or height. Skipping visual marker.")
                    # Fallback to default initial position for text box
                    initial_text_box_left = Inches(1.0).emu
                    initial_text_box_top = Inches(1.0).emu

                text_box_final_left, text_box_final_top = find_non_overlapping_position(
                    slide, initial_text_box_left, initial_text_box_top, text_box_collision_width, text_box_collision_height)

            except (ValueError, TypeError, KeyError, ZeroDivisionError) as e:
                logger.warning(
                    f"Could not calculate visual marker or text box position from rect for {annotation_id}: {annotation_rect}. Error: {e}. Using default position for text box.")
                # Fallback with default initial position for text box
                text_box_final_left, text_box_final_top = find_non_overlapping_position(
                    slide, Inches(1.0).emu, Inches(1.0).emu, text_box_collision_width, text_box_collision_height)
        else:
            if not largest_image:
                logger.warning(
                    f"No image found on slide for relative annotation positioning for {annotation_id}.")
            if not annotation_rect:
                logger.warning(
                    f"Missing or incomplete 'rect' data for annotation: {annotation_id}")
            logger.warning(
                f"Using default position for annotation text box {annotation_id}. Visual marker will not be drawn.")
            text_box_final_left, text_box_final_top = find_non_overlapping_position(
                slide, Inches(1.0).emu, Inches(1.0).emu, text_box_collision_width, text_box_collision_height)

        # Set current properties for the text box of an annotation
        current_shape_left = text_box_final_left
        current_shape_top = text_box_final_top
        current_shape_width = text_box_collision_width  # Will be adjusted by auto_size
        current_shape_height = text_box_collision_height  # Will be adjusted by auto_size
        current_shape_type = text_box_shape_type
        current_fill_color = text_box_fill_color
        current_line_color = text_box_line_color
        current_text_color = text_box_text_color
        current_shape_name = f"Annotation_{annotation_id}"

    # --- Add the Text Box / Comment Box Shape to the slide ---
    text_box_shape = None
    try:
        text_box_shape = slide.shapes.add_shape(
            current_shape_type, current_shape_left, current_shape_top, current_shape_width, current_shape_height
        )
        text_box_shape.name = current_shape_name

        # Apply Fill
        text_box_shape.fill.solid()
        text_box_shape.fill.fore_color.rgb = current_fill_color

        # Apply Line
        text_box_shape.line.color.rgb = current_line_color
        text_box_shape.line.width = Pt(1.5)

        # Add Text
        if annotation_text:
            text_frame = text_box_shape.text_frame

            main_text = annotation_text

            # Format the timestamp if available
            timestamp = get_annotation_value(
                annotation_data, 'timestamp', '')  # Use helper
            if timestamp:
                try:
                    # Format as YYYY/MM/DD if it's a valid timestamp
                    timestamp_str = timestamp.split('T')[0].replace('-', '/')
                except (IndexError, AttributeError):
                    timestamp_str = str(timestamp)
            else:
                timestamp_str = ""

            # Format start_at time if available for videos
            start_at = get_annotation_value(annotation_data, 'start_at')
            start_at_str = ""
            if start_at is not None and isinstance(start_at, (int, float)):
                try:
                    minutes = int(start_at // 60)
                    seconds = start_at % 60
                    # Format to MM:SS.ss and add play symbol
                    start_at_str = f"▶ {minutes:02d}:{seconds:05.2f}"
                except (TypeError, ValueError):
                    logger.warning(f"Could not format start_at: {start_at}")
                    start_at_str = ""

            # hard code for now.
            id_short = "Review FB"
            # Use retrieved ID for consistency if needed later
            # if annotation_id != 'unknown':
            #     id_short = annotation_id.split(
            #         '-')[0] if '-' in annotation_id else annotation_id[:8]

            # Create the header with date and ID
            header_parts = []
            if timestamp_str:
                header_parts.append(timestamp_str)
            if id_short:
                header_parts.append(id_short)

            header_text = "  ".join(header_parts)

            # Create the footer with additional metadata
            # footer_parts = []

            # footer_text = "  ".join(footer_parts)

            # Add left arrow if it's a comment
            if is_comment:
                main_text = f"← {main_text}"

            # Create the final text with header, main content, and footer
            if header_text and start_at_str:
                text_frame.text = f"{header_text}\n{start_at_str}\n{main_text}"
            elif header_text:
                text_frame.text = f"{header_text}\n{main_text}"
            elif start_at_str:
                text_frame.text = f"{start_at_str}\n{main_text}"
            else:
                text_frame.text = main_text

            # Text formatting
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            text_frame.margin_top = Inches(0.2)
            text_frame.margin_bottom = Inches(0.2)

            # Center text vertically for annotations, top for comments
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Format all paragraphs
            for i, paragraph in enumerate(text_frame.paragraphs):
                paragraph.alignment = PP_ALIGN.LEFT
                font = paragraph.font
                font.color.rgb = current_text_color  # Use current_text_color

                if i == 0:  # Header (date & ID) - medium size
                    font.size = Pt(10)
                    font.bold = True
                else:  # Main content - large size
                    font.size = Pt(12)
                    font.bold = True

        # --- Add Attachment Image if available ---
        attachment_url = get_annotation_value(
            annotation_data, 'attachment_image_url')  # Use helper
        s3_url = await get_s3_url_from_path(attachment_url)

        if s3_url and text_box_shape:  # Check if text_box_shape was successfully created
            logger.debug(
                f"Attempting to add attachment image from: {s3_url}")
            temp_file_path = None  # Initialize variable to store temp file path
            try:
                # Download image
                # Added log
                logger.info(f"Downloading image from {s3_url}")

                # Added timeout and redirects (use async client to avoid blocking loop)
                async with httpx.AsyncClient() as client:
                    response = await client.get(
                        s3_url, timeout=10, follow_redirects=True)
                response.raise_for_status()  # Raise exception for bad status codes

                content_type = response.headers.get('content-type', 'unknown')
                logger.info(f"Image content type: {content_type}")  # Added log
                # Basic check for image content type (optional but helpful)
                if not content_type.startswith('image/'):
                    logger.warning(
                        f"Downloaded content type ({content_type}) doesn't look like an image for {attachment_url}")
                    # Decide if you want to continue or skip based on this warning

                image_bytes = response.content
                if not image_bytes:
                    logger.warning(
                        f"Downloaded image content is empty for {attachment_url}. Skipping image add.")
                    return  # or continue, depending on desired behavior

                # --- Save image to temporary file ---
                # Try to guess extension, default to .png
                suffix = "." + \
                    content_type.split(
                        '/')[-1] if content_type.startswith('image/') else ".png"
                # File write can be heavy for large images; move to thread to avoid blocking

                def _save_temp(content: bytes, suffix: str) -> str:
                    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
                        temp_file.write(content)
                        return temp_file.name

                import asyncio
                temp_file_path = await asyncio.to_thread(_save_temp, image_bytes, suffix)
                logger.info(f"Image saved to temporary file: {temp_file_path}")

                # Position image below the annotation box
                img_left = text_box_shape.left
                # Add some padding (approx 0.05 inches)
                img_top = text_box_shape.top + \
                    text_box_shape.height + Emu(50000)

                # Add picture - let python-pptx determine initial size based on image DPI
                logger.info(
                    f"Adding picture to slide {slide.slide_id} at ({img_left}, {img_top}) for annotation {annotation_id} using path: {temp_file_path}")
                pic = slide.shapes.add_picture(
                    temp_file_path, img_left, img_top)
                # Added log
                logger.info(
                    f"Picture added initially with dimensions: Width={pic.width}, Height={pic.height}")

                # Resize picture to a reasonable maximum size, preserving aspect ratio
                max_dimension = Inches(2.0)  # Max width or height of 2 inches
                # Add check for positive dimensions before resizing
                if pic.width > 0 and pic.height > 0:
                    if pic.width > max_dimension or pic.height > max_dimension:
                        ratio = min(max_dimension / pic.width,
                                    max_dimension / pic.height)
                        pic.width = int(pic.width * ratio)
                        pic.height = int(pic.height * ratio)
                        # Re-adjust top position slightly if needed after resizing (less critical)
                        # pic.top = shape.top + shape.height + Emu(50000)
                        # Added log
                        logger.info(
                            f"Picture resized to: Width={pic.width}, Height={pic.height}")
                else:
                    logger.warning(
                        f"Initial picture dimensions were zero or negative. Skipping resize for annotation {annotation_id}")

                # Changed log message
                logger.info(
                    f"Successfully added and processed attachment image for annotation {annotation_id}")

            except httpx.RequestError as e:
                logger.error(
                    f"Network error downloading attachment image {attachment_url}: {e}")
            except httpx.HTTPStatusError as e:
                logger.error(
                    f"HTTP error downloading attachment image {attachment_url}: Status {e.response.status_code}")
            except Exception as e:
                # Catch potential errors from add_picture (e.g., invalid image format) or resizing
                # Use retrieved ID
                logger.error(
                    f"Failed to add or resize attachment image {attachment_url} for annotation {annotation_id}: {e}", exc_info=True)
            finally:
                # --- Clean up temporary file ---
                if temp_file_path and os.path.exists(temp_file_path):
                    try:
                        os.unlink(temp_file_path)
                        logger.info(
                            f"Cleaned up temporary image file: {temp_file_path}")
                    except OSError as e:
                        logger.error(
                            f"Error deleting temporary file {temp_file_path}: {e}")

    except Exception as e:
        logger.error(
            # Use retrieved ID, and current_shape_name for context
            f"Failed to add shape to slide for annotation/comment {annotation_id} (shape name: {current_shape_name}): {e}", exc_info=True)
