import json
from aimage_supervision.settings import logger
import os
import uuid
from datetime import UTC, datetime
from io import BytesIO
from typing import Any, Dict, List, Optional

import aioboto3
from dotenv import load_dotenv
from PIL import Image
from pptx import Presentation  # type: ignore

from aimage_supervision.clients.aws_s3 import (download_file_from_s3,
                                               upload_file_to_s3)
from aimage_supervision.utils.image_compression import compress_image_async


# 加载环境变量
load_dotenv()


class PPTXProcessor:
    def __init__(self):
        # 初始化S3会话
        self.session = aioboto3.Session(
            aws_access_key_id=os.getenv('AWS_ACCESS_KEY_ID'),
            aws_secret_access_key=os.getenv('AWS_SECRET_ACCESS_KEY'),
            region_name=os.getenv('AWS_REGION')
        )
        self.S3_BUCKET = os.getenv('AWS_BUCKET_NAME')

        # 图片处理配置
        self.IMAGE_CONFIG = {
            'max_width': 1200,  # 最大宽度
            'max_height': 800,  # 最大高度
            'quality': 85,      # JPEG压缩质量
            'optimize': True    # 优化图片
        }

    def _clean_text(self, text: str) -> str:
        """
        清理从 PPT 中提取的文本

        Args:
            text (str): 原始文本

        Returns:
            str: 清理后的文本
        """
        # 移除 PowerPoint 页码占位符
        text = text.replace('‹#›', '')

        # 清理多余的空白字符
        text = ' '.join(text.split())

        # 如果清理后文本为空，返回空字符串
        return text.strip() or ''

    def _process_image(self, image_blob: bytes) -> bytes:
        """
        处理图片：压缩和调整大小

        Args:
            image_blob (bytes): 原始图片二进制数据

        Returns:
            bytes: 处理后的图片二进制数据
        """
        try:
            # 从二进制数据加载图片
            img = Image.open(BytesIO(image_blob))

            # 获取原始尺寸
            original_width, original_height = img.size
            logger.info(f"原始图片尺寸: {original_width}x{original_height}")

            # 计算调整后的尺寸（等比缩放）
            max_width = self.IMAGE_CONFIG['max_width']
            max_height = self.IMAGE_CONFIG['max_height']

            # 只有当图片超过最大尺寸时才进行缩放
            if original_width > max_width or original_height > max_height:
                # 计算宽高比
                width_ratio = max_width / original_width
                height_ratio = max_height / original_height

                # 使用较小的比例进行缩放，以确保图片完全适应最大尺寸
                scale_ratio = min(width_ratio, height_ratio)

                new_width = int(original_width * scale_ratio)
                new_height = int(original_height * scale_ratio)

                # 调整图片大小
                img = img.resize((new_width, new_height), Image.LANCZOS)
                logger.info(f"调整后的图片尺寸: {new_width}x{new_height}")

            # 保存为优化的PNG或JPEG
            output = BytesIO()

            # 根据原始格式选择保存格式
            if img.format == 'PNG':
                img.save(output, format='PNG',
                         optimize=self.IMAGE_CONFIG['optimize'])
            else:
                # 如果不是PNG，转换为JPEG以获得更好的压缩效果
                if img.mode == 'RGBA':
                    # JPEG不支持透明度，需要转换
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    background.paste(img, mask=img.split()[3])  # 使用alpha通道作为蒙版
                    img = background

                img.save(output, format='JPEG',
                         quality=self.IMAGE_CONFIG['quality'],
                         optimize=self.IMAGE_CONFIG['optimize'])

            # 获取处理后的二进制数据
            processed_image = output.getvalue()
            logger.info(
                f"图片压缩: 原始大小 {len(image_blob)/1024:.2f}KB, 处理后 {len(processed_image)/1024:.2f}KB")

            return processed_image
        except Exception as e:
            logger.error(f"处理图片时出错: {str(e)}")
            # 如果处理失败，返回原始图片
            return image_blob

    async def process_pptx_file(self, pptx_file_path: str, output_folder: str, s3_folder: str) -> Dict[str, Any]:
        """
        处理单个PPT文件，提取内容并上传媒体文件到S3

        Args:
            pptx_file_path (str): PPT文件的本地路径
            output_folder (str): 输出文件夹路径
            s3_folder (str): S3目标文件夹路径

        Returns:
            Dict[str, Any]: 包含子任务列表和第一页文本的字典
        """
        logger.info(f"Opening presentation: {pptx_file_path}")
        presentation = Presentation(pptx_file_path)
        subtasks: List[Dict[str, Any]] = []
        first_slide_texts: List[str] = []

        # 获取幻灯片尺寸用于过滤图片
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height
        slide_area = slide_width * slide_height if slide_width and slide_height else 0

        filename_base = os.path.splitext(os.path.basename(pptx_file_path))[0]

        # 创建输出目录
        image_folder = os.path.join(output_folder, f"{filename_base}_images")
        os.makedirs(image_folder, exist_ok=True)

        # 创建媒体输出目录
        media_folder = os.path.join(output_folder, f"{filename_base}_media")
        os.makedirs(media_folder, exist_ok=True)

        # 提取第一页的文本（如果存在）
        for shape in presentation.slides[0].shapes:
            if shape.has_text_frame:
                cleaned_text = self._clean_text(shape.text)
                if cleaned_text:
                    first_slide_texts.append(cleaned_text)

        for slide_number, slide in enumerate(presentation.slides, 1):
            logger.info(f"Processing slide {slide_number}")

            # 提取文本
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    cleaned_text = self._clean_text(shape.text)
                    if cleaned_text:
                        texts.append(cleaned_text)

            # --- 新逻辑：优先处理视频，采用"左上优先"（带误差范围）原则选择主资产 ---
            has_video = False
            TOLERANCE_RATIO = 0.15  # 15%的误差容忍度
            margin = slide_width * TOLERANCE_RATIO if slide_width else 0

            video_candidates = [
                s for s in slide.shapes if s.shape_type == 16]

            if video_candidates:
                best_video_candidate = None
                min_left, min_top = float('inf'), float('inf')

                for candidate in video_candidates:
                    if not (hasattr(candidate, 'left') and hasattr(candidate, 'top')):
                        continue

                    # 判断逻辑：
                    # 1. 如果新候选者在当前最佳者的左侧（超出误差范围），则更新。
                    # 2. 如果在同一垂直带内（在误差范围内），则选择更靠上的一个。
                    is_in_new_left_column = candidate.left < (
                        min_left - margin)
                    is_in_same_column_and_higher = abs(
                        candidate.left - min_left) <= margin and candidate.top < min_top

                    if is_in_new_left_column or is_in_same_column_and_higher:
                        min_left, min_top = candidate.left, candidate.top
                        best_video_candidate = candidate

                if best_video_candidate:
                    try:
                        logger.info(
                            f"Processing main video asset on slide {slide_number}")
                        has_video = True

                        # 使用现有的提取逻辑，它会获取幻灯片关系中的第一个媒体文件。
                        # 这是一个已知的局限，如果幻灯片有多个视频文件，可能无法精确匹配。
                        rels = slide.part.rels
                        for rel in rels.values():
                            if 'media' in rel.reltype:
                                media_part = rel.target_part
                                media_filename = f"{filename_base}_slide{slide_number}.mp4"
                                media_path = os.path.join(
                                    media_folder, media_filename)

                                # 保存到本地
                                with open(media_path, "wb") as media_file:
                                    media_file.write(media_part.blob)

                                # 上传到S3
                                s3_path = f"{s3_folder}/media/{media_filename}"
                                with open(media_path, 'rb') as media_file:
                                    await upload_file_to_s3(media_file, s3_path)

                                # 使用清理后的文本
                                title = texts[0] if texts else f"Slide {slide_number}"

                                subtasks.append({
                                    "title": title,
                                    "s3_path": s3_path,
                                    "description": "\n".join(texts),
                                    "task_type": "video",
                                    "marked": [],
                                    "slide_page_number": slide_number
                                })
                                # 只处理找到的第一个媒体关系
                                break
                    except Exception as e:
                        logger.error(
                            f"Error processing video in slide {slide_number}: {str(e)}")

            # 如果已经处理了视频，则跳过处理图片
            if has_video:
                continue

            # --- 新逻辑：基于"左上优先"（带误差范围）位置原则选择唯一的资产图片 ---
            main_asset_candidate: Optional[Any] = None
            min_left, min_top = float('inf'), float('inf')

            for shape in slide.shapes:
                if shape.shape_type == 13:  # 图片
                    # Heuristic: filter out small images that are likely logos or decorative elements.
                    if slide_area > 0 and (shape.width * shape.height / slide_area) < 0.03:
                        logger.info(
                            f"Skipping small image on slide {slide_number} - likely a logo or decorative element."
                        )
                        continue

                    if not (hasattr(shape, 'left') and hasattr(shape, 'top')):
                        continue

                    # 判断逻辑（同视频选择逻辑）：
                    # 1. 如果新候选者在当前最佳者的左侧（超出误差范围），则更新。
                    # 2. 如果在同一垂直带内（在误差范围内），则选择更靠上的一个。
                    is_in_new_left_column = shape.left < (min_left - margin)
                    is_in_same_column_and_higher = abs(
                        shape.left - min_left) <= margin and shape.top < min_top

                    if is_in_new_left_column or is_in_same_column_and_higher:
                        min_left, min_top = shape.left, shape.top
                        main_asset_candidate = shape

            # 如果找到了主要资产，则只处理这张图片并创建subtask
            if main_asset_candidate:
                try:
                    image = main_asset_candidate.image
                    image_blob = image.blob
                    image_format = image.content_type.split('/')[-1].lower()

                    # 优先使用png,jpeg等常见格式
                    if image_format not in ['png', 'jpeg', 'jpg', 'gif']:
                        image_format = 'png'

                    # 使用shape_id确保文件名唯一
                    base_image_filename = f"{filename_base}_slide{slide_number}_{main_asset_candidate.shape_id}"
                    original_image_filename = f"{base_image_filename}.{image_format}"

                    # 1. 上传原始图片
                    original_s3_path = f"{s3_folder}/images/original/{original_image_filename}"
                    with BytesIO(image_blob) as img_file:
                        await upload_file_to_s3(img_file, original_s3_path)

                    # 2. 压缩图片并上传
                    compressed_s3_path = None
                    try:
                        with BytesIO(image_blob) as img_file:
                            compressed_file, output_format, _ = await compress_image_async(img_file)

                        compressed_ext = 'jpg' if output_format == 'JPEG' else output_format.lower()
                        compressed_image_filename = f"{base_image_filename}_compressed.{compressed_ext}"

                        compressed_s3_path = f"{s3_folder}/images/compressed/{compressed_image_filename}"

                        # Check if compressed_file is the same as original (no compression benefit)
                        if compressed_file.closed or hasattr(compressed_file, 'name'):
                            # If original file was returned (no compression benefit), create a new BytesIO
                            compressed_file = BytesIO(image_blob)

                        await upload_file_to_s3(compressed_file, compressed_s3_path)
                    except Exception as e:
                        logger.error(
                            f"Error compressing image on slide {slide_number}: {str(e)}")
                        # 压缩失败不影响主流程，只记录日志

                    # 使用清理后的文本
                    title = texts[0] if texts else f"Slide {slide_number}"

                    subtasks.append({
                        "title": title,
                        "s3_path": original_s3_path,
                        "compressed_s3_path": compressed_s3_path,
                        "description": "\n".join(texts),
                        "task_type": "picture",
                        "marked": [],
                        "slide_page_number": slide_number
                    })
                except Exception as e:
                    logger.error(
                        f"Error processing main asset image on slide {slide_number}: {str(e)}")

        return {
            'subtasks': subtasks,
            'first_slide_description': "\n".join(first_slide_texts) if first_slide_texts else filename_base,
        }
