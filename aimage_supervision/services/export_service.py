import asyncio
import io
import json
from aimage_supervision.settings import logger
import os
import re
import tempfile
from typing import Dict, List, Union

import httpx
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.util import Emu, Inches, Pt

from aimage_supervision.clients.aws_s3 import (download_file_content_from_s3,
                                               get_s3_url_from_path)
from aimage_supervision.enums import SubtaskType
from aimage_supervision.exceptions import (ContentNotFoundError,
                                           FileTooLargeError, InvalidPathError)
from aimage_supervision.models import Subtask, Task
from aimage_supervision.utils.export_pptx import (add_comment_box_to_slide,
                                                  find_largest_image,
                                                  get_annotation_value,
                                                  hex_to_rgb)
from aimage_supervision.utils.file_converter import convert_pptx_to_pdf
from aimage_supervision.utils.video_utils import \
    get_frame_at_timestamp_from_path


class PptxExportService:
    """Unified service for exporting tasks to PPTX format"""

    def __init__(self):
        """Initialize the export service"""
        self.logger = logger

    async def export_task_to_pdf(self, task: Task) -> io.BytesIO:
        """Exports the task to a PDF by converting the generated PPTX."""
        self.logger.info(f"Starting PDF export for task {task.id}...")

        # 1. Reuse existing PPTX logic (DRY)
        pptx_buffer = await self.export_task_to_pptx(task)

        # 2. Convert to PDF using the new utility (KISS)
        pdf_bytes = convert_pptx_to_pdf(pptx_buffer.getvalue())

        self.logger.info(f"Successfully generated PDF for task {task.id}")
        return io.BytesIO(pdf_bytes)

    async def export_task_to_pptx(self, task: Task) -> io.BytesIO:
        """
        Main export method - handles both PPTX and image-based tasks

        Args:
            task: Task object with prefetched subtasks

        Returns:
            BytesIO: PPTX file content as bytes

        Raises:
            FileTooLargeError: If file size exceeds 300MB limit
            ContentNotFoundError: If no exportable content found
            InvalidPathError: If S3 path is invalid or empty
            Exception: For other export failures
        """
        self.logger.info(f"Starting PPTX export for task {task.id}")

        # Decision logic: Check if s3_path points to an actual PPTX file
        if task.s3_path and task.s3_path.strip() and task.s3_path.lower().endswith('.pptx'):
            self.logger.info(f"Exporting PPTX-based task {task.id}")
            return await self._export_from_existing_pptx(task)
        else:
            self.logger.info(f"Exporting media-based task {task.id}")
            return await self._export_from_media(task)

    async def _export_from_existing_pptx(self, task: Task) -> io.BytesIO:
        """
        Handle tasks created from PPTX files (existing logic)

        Args:
            task: Task with existing PPTX s3_path

        Returns:
            BytesIO: Modified PPTX with annotations

        Raises:
            ValueError: If S3 path is invalid
            Exception: For download or processing failures
        """
        s3_path = task.s3_path
        if not s3_path or not s3_path.strip():
            raise InvalidPathError("No valid S3 path found for PPTX export")

        self.logger.info(f"Exporting existing PPTX from S3 path: {s3_path}")

        # 下载PPTX文件 (DRY - reuse existing S3 download function)
        try:
            pptx_content = await download_file_content_from_s3(task.s3_path)
            self.logger.info(
                f"Downloaded PPTX file, size: {len(pptx_content)} bytes")
            if len(pptx_content) > 300 * 1024 * 1024:
                raise FileTooLargeError('PPTXファイルのサイズが300MBを超えているため、エクスポートできません（現在のサイズ: {}MB）'.format(
                    round(len(pptx_content) / (1024 * 1024), 2)))
        except FileTooLargeError:
            # Re-raise FileTooLargeError without wrapping
            raise
        except Exception as e:
            self.logger.error(f"Failed to download PPTX from S3: {e}")
            raise Exception(f"Failed to download PPTX file: {e}")
        try:
            pptx_buffer = io.BytesIO(pptx_content)
            # 构建 Presentation 可能较重，放到线程池中
            prs = await asyncio.to_thread(Presentation, pptx_buffer)

            # Process annotations
            subtasks = task.subtasks
            if not subtasks:
                self.logger.warning(
                    f"Task {task.id} has no subtasks, returning original PPTX.")
            else:
                # Get slide dimensions from the presentation object
                slide_width = prs.slide_width
                slide_height = prs.slide_height

                annotation_count = 0
                for subtask in subtasks:
                    if not subtask.annotations or not isinstance(subtask.annotations, list):
                        continue

                    # Convert from 1-based to 0-based index for prs.slides access
                    page_index = subtask.slide_page_number - 1
                    # Add bounds check for page_index
                    if 0 <= page_index < len(prs.slides):
                        slide = prs.slides[page_index]
                    else:
                        self.logger.warning(
                            f"Invalid page index {subtask.slide_page_number} (converted to {page_index}) for subtask {subtask.name} (Task {task.id}). Skipping annotations for this subtask.")
                        continue  # Skip this subtask if page index is invalid

                    # 用新的右上角布局方式处理所有注释
                    try:
                        annotation_count = await self._add_annotations_in_right_panel(
                            slide, subtask.annotations, slide_width, slide_height)
                        self.logger.info(
                            f"Added {annotation_count} annotations in right panel layout for slide")
                    except Exception as e:
                        self.logger.warning(
                            f"Failed to add annotations in right panel: {e}")
                        # Continue with other annotations or slide processing
                        continue

                self.logger.info(
                    f"Added {annotation_count} annotations to PPTX for task {task.id}")

            # Save modified presentation to a new buffer
            output_buffer = io.BytesIO()
            # 保存 PPTX 到内存可能较重，放到线程池中
            await asyncio.to_thread(prs.save, output_buffer)
            output_buffer.seek(0)

            self.logger.info(
                f"Successfully processed existing PPTX for task {task.id}")
            return output_buffer

        except (FileTooLargeError, ContentNotFoundError, InvalidPathError):
            # Re-raise our custom exceptions without wrapping
            raise
        except Exception as e:
            self.logger.error(
                f"Error processing PPTX for task {task.id}: {e}", exc_info=True)
            raise Exception(f"Failed to process PPTX file: {e}")

    async def _export_from_media(self, task: Task) -> io.BytesIO:
        """
        Handle tasks created from media (new logic)

        Args:
            task: Task with media-based subtasks

        Returns:
            BytesIO: Generated PPTX with media and annotations

        Raises:
            ValueError: If no media subtasks found
        """
        # Get media subtasks only
        media_subtasks = [s for s in task.subtasks if (
            s.task_type == SubtaskType.PICTURE or s.task_type == SubtaskType.VIDEO) and s.content]

        if not media_subtasks:
            raise ContentNotFoundError("No media subtasks found for export")

        self.logger.info(
            f"Found {len(media_subtasks)} media subtasks for task {task.id}")

        # 1. Create blank presentation
        prs = await self._create_blank_presentation(task)

        # 2. Add slide for each media subtask
        for subtask in media_subtasks:
            try:
                if subtask.task_type == SubtaskType.PICTURE:
                    await self._add_image_slide(prs, subtask)
                    self.logger.info(
                        f"Added slide for image subtask {subtask.id}")
                elif subtask.task_type == SubtaskType.VIDEO:
                    await self._add_video_slides(prs, subtask)
                    self.logger.info(
                        f"Added slides for video subtask {subtask.id}")
            except Exception as e:
                self.logger.warning(
                    f"Failed to add slide(s) for subtask {subtask.id}: {e}")
                # Continue with other subtasks
                continue

        # 3. Return as BytesIO
        output_buffer = io.BytesIO()
        prs.save(output_buffer)
        output_buffer.seek(0)

        self.logger.info(f"Successfully generated PPTX for task {task.id}")
        return output_buffer

    async def _create_blank_presentation(self, task: Task) -> PresentationType:
        """
        Create a new presentation with title slide

        Args:
            task: Task object for title information

        Returns:
            Presentation: New PPTX presentation with title slide
        """
        prs = Presentation()

        # Add title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        title_slide = prs.slides.add_slide(title_slide_layout)

        # Set title and subtitle
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]

        title.text = task.name or "Task Export"

        # Add project name or description as subtitle if available
        subtitle_text = ""
        if hasattr(task, 'project') and task.project:
            # Properly await the project relationship
            project = await task.project
            subtitle_text = f"Project: {project.name}"
        if task.description:
            if subtitle_text:
                subtitle_text += f"\n{task.description}"
            else:
                subtitle_text = task.description

        if subtitle_text:
            subtitle.text = subtitle_text

        self.logger.info(f"Created title slide for task {task.name}")
        return prs

    async def _add_image_slide(self, prs: PresentationType, subtask: Subtask) -> None:
        """
        Add a slide with media and annotations

        Args:
            prs: Presentation object
            subtask: Subtask containing media content

        Raises:
            Exception: If media cannot be added or processed
        """
        # 1. Add blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # 2. Add title if subtask has name
        if subtask.name:
            self._add_slide_title(slide, subtask.name)

        # 3. Download and add media
        try:
            await self._add_image_to_slide(slide, subtask)
            self.logger.info(f"Added media to slide for subtask {subtask.id}")
        except Exception as e:
            self.logger.warning(
                f"Failed to add media for subtask {subtask.id}: {e}")
            # Add error placeholder instead
            self._add_error_placeholder(
                slide, f"Media unavailable: {subtask.name}")

        # 4. Add annotations using new right panel layout
        if subtask.annotations:
            self.logger.info(
                f"Adding {len(subtask.annotations)} annotations for subtask {subtask.id}")
            try:
                annotation_count = await self._add_annotations_in_right_panel(
                    slide, subtask.annotations, prs.slide_width, prs.slide_height)
                self.logger.info(
                    f"Added {annotation_count} annotations in right panel layout")
            except Exception as e:
                self.logger.warning(
                    f"Failed to add annotations in right panel: {e}")

    async def _add_video_slides(self, prs: PresentationType, subtask: Subtask):
        """
        Add multiple slides for a video subtask based on annotations with timestamps.

        Args:
            prs (Presentation): The presentation object to add slides to.
            subtask (Subtask): The video subtask with annotations.
        """
        if not subtask.annotations:
            self.logger.info(
                f"Video subtask {subtask.id} has no annotations to export.")
            return

        s3_path = self._get_s3_path_from_subtask(subtask)
        if not s3_path:
            self.logger.warning(
                f"No S3 path found for video subtask {subtask.id}, skipping.")
            # Optionally add a slide indicating the error
            self._add_full_slide_error_placeholder(
                prs, subtask.name or f"Subtask {subtask.id}", "Video content is missing or path is invalid.")
            return

        # Download video to a temporary file once
        video_temp_path = None
        try:
            video_bytes = await download_file_content_from_s3(s3_path)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".mp4") as temp_file:
                temp_file.write(video_bytes)
                video_temp_path = temp_file.name
            self.logger.info(
                f"Video for subtask {subtask.id} downloaded to {video_temp_path}")

            # Group annotations by timestamp to handle multiple comments on the same frame
            annotations_by_ts: Dict[float, List] = {}
            for annotation in subtask.annotations:
                # Per user feedback, "start_at" is the correct field for the video timestamp, not "timestamp".
                timestamp = getattr(annotation, 'start_at', None)
                # A timestamp of 0.0 is valid, so we must explicitly check for None.
                if timestamp is None:
                    # Skip annotations that are not tied to a specific time in the video.
                    continue

                try:
                    # The 'start_at' value should be a float or int representing seconds.
                    ts_key = float(timestamp)
                except (ValueError, TypeError) as e:
                    self.logger.warning(
                        f"Skipping annotation with invalid 'start_at' format for subtask {subtask.id}. "
                        f"Expected a number, but got: '{timestamp}'. Error: {e}"
                    )
                    continue

                if ts_key not in annotations_by_ts:
                    annotations_by_ts[ts_key] = []
                annotations_by_ts[ts_key].append(annotation)

            if not annotations_by_ts:
                self.logger.warning(
                    f"No valid time-based annotations found for video subtask {subtask.id}.")
                return

            # Process each timestamped frame
            for timestamp, annotations in sorted(annotations_by_ts.items()):
                self.logger.info(
                    f"Processing frame for subtask {subtask.id} at timestamp {timestamp}s")
                try:
                    # 抽帧使用 moviepy，同步且CPU密集，放线程池
                    frame_bytes = await asyncio.to_thread(
                        get_frame_at_timestamp_from_path, video_temp_path, float(
                            timestamp)
                    )

                    # Create a new slide for this frame
                    slide = prs.slides.add_slide(
                        prs.slide_layouts[6])  # Blank layout

                    # Add title with timestamp
                    title = f"{subtask.name} (Timestamp: {timestamp}s)"
                    self._add_slide_title(slide, title)

                    # Add frame image to slide
                    pic = self._add_image_bytes_to_slide(slide, frame_bytes)
                    self._scale_image_to_fit(pic, Inches(
                        8), Inches(5.5))  # Reuse scaling logic

                    # Add all annotations for this timestamp
                    for ann in annotations:
                        await add_comment_box_to_slide(slide, ann, prs.slide_width, prs.slide_height)

                except Exception as e:
                    self.logger.error(
                        f"Failed to process frame at timestamp {timestamp} for subtask {subtask.id}: {e}")
                    # Optionally add an error slide for this specific frame
                    self._add_full_slide_error_placeholder(
                        prs, subtask.name or f"Subtask {subtask.id}", f"Failed to extract frame at {timestamp}s: {e}")
                    continue

        except Exception as e:
            self.logger.error(
                f"Failed to process video subtask {subtask.id}: {e}", exc_info=True)
            self._add_full_slide_error_placeholder(
                prs, subtask.name or f"Subtask {subtask.id}", f"Could not process video: {e}")
        finally:
            # Clean up the temporary video file
            if video_temp_path and os.path.exists(video_temp_path):
                try:
                    os.unlink(video_temp_path)
                    self.logger.info(
                        f"Cleaned up temp video file: {video_temp_path}")
                except OSError as e:
                    self.logger.warning(
                        f"Failed to delete temp video file {video_temp_path}: {e}")

    def _add_slide_title(self, slide, title_text: str) -> None:
        """
        Add title to a slide

        Args:
            slide: Slide object
            title_text: Title text to add
        """
        # Add title text box at the top of the slide
        left = Inches(0.5)
        top = Inches(0.2)
        width = Inches(9)
        height = Inches(0.8)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title_text

        # Format title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(24)
        title_paragraph.font.bold = True

    def _get_s3_path_from_subtask(self, subtask: Subtask) -> Union[str, None]:
        """Extracts S3 path from various possible subtask.content structures."""
        if not subtask.content:
            return None

        # Handle Pydantic model, dict, and other cases
        if hasattr(subtask.content, 's3_path'):
            return subtask.content.s3_path
        elif isinstance(subtask.content, dict):
            return subtask.content.get('s3_path')
        else:
            try:
                if hasattr(subtask.content, 'model_dump'):
                    content_dict = subtask.content.model_dump()
                elif hasattr(subtask.content, 'dict'):
                    content_dict = subtask.content.dict()
                else:
                    return None
                return content_dict.get('s3_path')
            except Exception:
                return None

    def _add_image_bytes_to_slide(self, slide, image_bytes: bytes):
        """Adds image from bytes to a slide and returns the shape."""
        # Use a temporary file to add the picture from bytes
        with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as temp_image_file:
            temp_image_file.write(image_bytes)
            temp_image_path = temp_image_file.name

        try:
            left = Inches(1)
            top = Inches(1.5)  # Leave space for title
            pic = slide.shapes.add_picture(temp_image_path, left, top)
            return pic
        finally:
            # Clean up the temp image file
            if os.path.exists(temp_image_path):
                os.unlink(temp_image_path)

    async def _add_image_to_slide(self, slide, subtask: Subtask):
        """
        Download media and add to slide with proper sizing

        Args:
            slide: Slide object
            subtask: Subtask containing media content

        Returns:
            Shape: Added media shape

        Raises:
            Exception: If media cannot be downloaded or added
        """
        # 1. Get S3 path from subtask content (handle both dict and Pydantic model)
        self.logger.info(
            f"Processing subtask {subtask.id} with content type: {type(subtask.content)}")

        s3_path = self._get_s3_path_from_subtask(subtask)

        if not s3_path or not s3_path.strip():
            raise InvalidPathError(
                f"Empty S3 path in subtask {subtask.id} content")

        self.logger.info(f"Processing media from S3 path: {s3_path}")

        # 2. Download image directly from S3 (DRY - reuse existing function)
        try:
            media_bytes = await download_file_content_from_s3(s3_path)
            self.logger.info(f"Downloaded media: {len(media_bytes)} bytes")
        except InvalidPathError:
            # Re-raise InvalidPathError without wrapping
            raise
        except Exception as e:
            self.logger.error(f"Failed to download media from S3: {e}")
            raise Exception(f"Failed to download media: {e}")

        # 3. Add to slide (centered, fit to content area)
        pic = self._add_image_bytes_to_slide(slide, media_bytes)

        # 4. Scale to fit while preserving aspect ratio
        max_width = Inches(8)
        max_height = Inches(5.5)
        self._scale_image_to_fit(pic, max_width, max_height)

        self.logger.info(
            f"Successfully added media to slide for subtask {subtask.id}")
        return pic

    def _scale_image_to_fit(self, image_shape, max_width, max_height) -> None:
        """
        Scale image to fit within max dimensions while preserving aspect ratio

        Args:
            image_shape: Image shape object
            max_width: Maximum width
            max_height: Maximum height
        """
        # Calculate scale factors
        width_scale = max_width / image_shape.width
        height_scale = max_height / image_shape.height

        # Use the smaller scale to maintain aspect ratio
        scale = min(width_scale, height_scale)

        # Only scale down, not up
        if scale < 1:
            image_shape.width = int(image_shape.width * scale)
            image_shape.height = int(image_shape.height * scale)

    def _add_error_placeholder(self, slide, error_message: str) -> None:
        """
        Add error placeholder text when image cannot be loaded

        Args:
            slide: Slide object
            error_message: Error message to display
        """
        # Add error message text box
        left = Inches(2)
        top = Inches(3)
        width = Inches(6)
        height = Inches(2)

        error_box = slide.shapes.add_textbox(left, top, width, height)
        error_frame = error_box.text_frame
        error_frame.text = error_message

        # Format error text
        error_paragraph = error_frame.paragraphs[0]
        error_paragraph.font.size = Pt(16)
        error_paragraph.font.italic = True

        self.logger.info(f"Added error placeholder: {error_message}")

    def _add_full_slide_error_placeholder(self, prs: PresentationType, title: str, message: str):
        """Adds a full slide dedicated to showing an error message."""
        slide = prs.slides.add_slide(
            prs.slide_layouts[1])  # Use layout 1: Title and Content
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]

        title_shape.text = f"Error: {title}"

        text_frame = body_shape.text_frame
        text_frame.text = message
        p = text_frame.paragraphs[0]
        p.font.size = Pt(18)

        self.logger.error(f"Added error slide for '{title}': {message}")

    async def _add_annotations_in_right_panel(self, slide, annotations, slide_width, slide_height):
        """
        将所有注释统一放在右上角区域管理，保留原图上的红框标记

        Args:
            slide: PPT幻灯片对象
            annotations: 注释列表
            slide_width: 幻灯片宽度
            slide_height: 幻灯片高度

        Returns:
            int: 成功添加的注释数量
        """
        if not annotations:
            return 0

        # 过滤掉AI注释
        filtered_annotations = [
            ann for ann in annotations
            if getattr(ann, 'type', None) not in ['ai-comment', 'ai-annotation']
        ]

        if not filtered_annotations:
            return 0

        self.logger.info(
            f"Processing {len(filtered_annotations)} annotations for right panel layout")

        # 找到幻灯片上最大的图像
        largest_image = find_largest_image(slide, slide_width, slide_height)

        annotation_count = 0

        # 1. 处理绘画数据 - 合成到原图上
        for annotation in filtered_annotations:
            if hasattr(annotation, 'drawing_data') or (isinstance(annotation, dict) and annotation.get('drawing_data')):
                await self._apply_drawing_to_image(slide, annotation, largest_image)

        # 2. 首先添加视觉标记框（保留原图上的红框）
        for i, annotation in enumerate(filtered_annotations, 1):
            try:
                await self._add_visual_marker_only(slide, annotation, largest_image, i)
                self.logger.info(
                    f"Added visual marker {i} for annotation {get_annotation_value(annotation, 'id', 'unknown')}")
            except Exception as e:
                self.logger.warning(
                    f"Failed to add visual marker for annotation {i}: {e}")
                continue

        # 3. 然后在右上角区域添加注释面板
        try:
            panel_count = await self._add_right_panel_annotations(
                slide, filtered_annotations, slide_width, slide_height)
            annotation_count = panel_count
            self.logger.info(f"Added {panel_count} annotations to right panel")
        except Exception as e:
            self.logger.warning(f"Failed to add right panel annotations: {e}")

        return annotation_count

    async def _add_visual_marker_only(self, slide, annotation, largest_image, marker_number):
        """
        仅添加视觉标记框到原图上，不添加文本框

        Args:
            slide: PPT幻灯片对象  
            annotation: 注释对象
            largest_image: 最大图像对象
            marker_number: 标记序号
        """
        annotation_rect = get_annotation_value(annotation, 'rect', None)
        if not annotation_rect or not largest_image:
            return

        # 获取注释的颜色
        annotation_color_hex = get_annotation_value(
            annotation, 'color', '#ff0000')  # 默认红色
        try:
            marker_color = hex_to_rgb(
                annotation_color_hex) if annotation_color_hex else RGBColor(255, 0, 0)
        except (ValueError, TypeError):
            marker_color = RGBColor(255, 0, 0)  # 默认红色

        # 计算DPI和位置（复用原有逻辑）
        img_left_emu = largest_image.left
        img_top_emu = largest_image.top
        current_image_width_emu = largest_image.width
        current_image_height_emu = largest_image.height

        # DPI计算逻辑（复用原有代码）
        BASELINE_DPI = 220.0
        BASELINE_WIDTH_EMU = 8000000.0
        BASELINE_HEIGHT_EMU = 3400000.0

        assumed_dpi = BASELINE_DPI
        if current_image_width_emu > 0 and current_image_height_emu > 0:
            is_landscape = (current_image_width_emu /
                            current_image_height_emu) > 1.0

            if is_landscape:
                dynamic_dpi = BASELINE_DPI * \
                    (BASELINE_WIDTH_EMU / current_image_width_emu)
                assumed_dpi = dynamic_dpi
            else:
                dynamic_dpi = BASELINE_DPI * \
                    (BASELINE_HEIGHT_EMU / current_image_height_emu)
                assumed_dpi = dynamic_dpi

        # 计算标记框位置
        visual_rect_x_pixels = annotation_rect.get('x', 0) if hasattr(
            annotation_rect, 'get') else getattr(annotation_rect, 'x', 0)
        visual_rect_y_pixels = annotation_rect.get('y', 0) if hasattr(
            annotation_rect, 'get') else getattr(annotation_rect, 'y', 0)
        visual_rect_width_pixels = annotation_rect.get('width', 0) if hasattr(
            annotation_rect, 'get') else getattr(annotation_rect, 'width', 0)
        visual_rect_height_pixels = annotation_rect.get('height', 0) if hasattr(
            annotation_rect, 'get') else getattr(annotation_rect, 'height', 0)

        visual_rect_x_emu = Inches(visual_rect_x_pixels / assumed_dpi).emu
        visual_rect_y_emu = Inches(visual_rect_y_pixels / assumed_dpi).emu
        visual_rect_width_emu = Inches(
            visual_rect_width_pixels / assumed_dpi).emu
        visual_rect_height_emu = Inches(
            visual_rect_height_pixels / assumed_dpi).emu

        if visual_rect_width_emu > 0 and visual_rect_height_emu > 0:
            visual_shape_actual_left = Emu(img_left_emu + visual_rect_x_emu)
            visual_shape_actual_top = Emu(img_top_emu + visual_rect_y_emu)

            # 添加标记框
            tool = get_annotation_value(annotation, 'tool', 'rect')
            shape_type = MSO_SHAPE.OVAL if tool == 'ellipse' else MSO_SHAPE.RECTANGLE

            marker_shape = slide.shapes.add_shape(
                shape_type,
                visual_shape_actual_left,
                visual_shape_actual_top,
                Emu(visual_rect_width_emu),
                Emu(visual_rect_height_emu)
            )

            marker_shape.name = f"VisualMarker_{marker_number}"
            marker_shape.fill.background()  # 透明填充
            marker_shape.line.color.rgb = marker_color
            marker_shape.line.width = Pt(2.0)

            # 在标记框中心添加序号
            # 计算序号文本框的位置（标记框中心）
            number_left = visual_shape_actual_left
            number_top = visual_shape_actual_top
            number_width = Emu(visual_rect_width_emu)
            number_height = Emu(visual_rect_height_emu)

            number_shape = slide.shapes.add_textbox(
                number_left, number_top, number_width, number_height
            )
            number_shape.name = f"MarkerNumber_{marker_number}"

            # 设置序号文本
            text_frame = number_shape.text_frame
            text_frame.text = str(marker_number)
            text_frame.word_wrap = False

            # 设置文本格式
            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = 1  # 居中对齐 (PP_ALIGN.CENTER)
            text_frame.vertical_anchor = 1  # 垂直居中 (MSO_ANCHOR.MIDDLE)

            font = paragraph.font
            font.size = Pt(14)
            font.bold = True
            font.color.rgb = marker_color

            # 去除文本框边框和填充
            number_shape.fill.background()
            number_shape.line.fill.background()

    async def _add_right_panel_annotations(self, slide, annotations, slide_width, _slide_height):
        """
        在右上角区域添加注释面板

        Args:
            slide: PPT幻灯片对象
            annotations: 注释列表  
            slide_width: 幻灯片宽度
            slide_height: 幻灯片高度

        Returns:
            int: 成功添加的注释数量
        """
        # 右上角面板配置
        panel_width = Inches(3.5)  # 面板宽度
        panel_start_left = slide_width - \
            panel_width - Inches(0.3)  # 距离右边缘0.3英寸
        panel_start_top = Inches(0.5)  # 距离顶部0.5英寸

        current_top = panel_start_top
        annotation_count = 0

        for i, annotation in enumerate(annotations, 1):
            try:
                # 创建注释文本框
                annotation_text = str(
                    get_annotation_value(annotation, 'text', ''))
                if not annotation_text:
                    continue

                # 计算文本框高度（基于文本长度的估算）
                estimated_lines = max(
                    1, len(annotation_text) // 40 + 1)  # 假设每行40字符
                text_height = Inches(0.3 + estimated_lines * 0.2)  # 基础高度+行高

                # 创建注释框
                annotation_box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    panel_start_left,
                    current_top,
                    panel_width,
                    text_height
                )
                annotation_box.name = f"RightPanelAnnotation_{i}"

                # 设置外观
                annotation_box.fill.solid()
                annotation_box.fill.fore_color.rgb = RGBColor(
                    239, 157, 83)  # 橙色背景
                annotation_box.line.color.rgb = RGBColor(150, 100, 50)  # 深橙色边框
                annotation_box.line.width = Pt(1.5)

                # 设置文本内容
                text_frame = annotation_box.text_frame
                text_frame.margin_left = Inches(0.15)
                text_frame.margin_right = Inches(0.15)
                text_frame.margin_top = Inches(0.1)
                text_frame.margin_bottom = Inches(0.1)
                text_frame.word_wrap = True

                # 格式化文本内容
                timestamp = get_annotation_value(annotation, 'timestamp', '')
                if timestamp:
                    try:
                        timestamp_str = timestamp.split(
                            'T')[0].replace('-', '/')
                    except (IndexError, AttributeError):
                        timestamp_str = str(timestamp)
                else:
                    timestamp_str = ''

                # 组装文本：序号 + 时间戳 + 内容
                header_parts = [f"{i}"]
                if timestamp_str:
                    header_parts.append("Review FB")

                header_text = "  ".join(header_parts)

                if timestamp_str:
                    full_text = f"{timestamp_str}\n{header_text}\n{annotation_text}"
                else:
                    full_text = f"{header_text}\n{annotation_text}"

                text_frame.text = full_text

                # 设置文本格式
                for j, paragraph in enumerate(text_frame.paragraphs):
                    paragraph.alignment = PP_ALIGN.LEFT
                    font = paragraph.font
                    font.color.rgb = RGBColor(255, 255, 255)  # 白色文字

                    if j == 0:  # 时间戳行
                        font.size = Pt(9)
                        font.bold = True
                    elif j == 1:  # 序号和标题行
                        font.size = Pt(10)
                        font.bold = True
                    else:  # 内容行
                        font.size = Pt(11)
                        font.bold = False

                current_annotation_bottom = current_top + text_height

                # 处理附件图片
                attachment_url = get_annotation_value(
                    annotation, 'attachment_image_url', None)
                if attachment_url:
                    try:
                        # 获取S3 URL
                        s3_url = await get_s3_url_from_path(attachment_url)
                        if s3_url:
                            # 下载图片
                            async with httpx.AsyncClient() as client:
                                response = await client.get(
                                    s3_url, timeout=10, follow_redirects=True)
                            response.raise_for_status()

                            image_bytes = response.content
                            if image_bytes:
                                # 保存到临时文件
                                content_type = response.headers.get(
                                    'content-type', 'unknown')
                                suffix = "." + \
                                    content_type.split(
                                        '/')[-1] if content_type.startswith('image/') else ".png"

                                with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
                                    temp_file.write(image_bytes)
                                    temp_file_path = temp_file.name

                                # 添加图片到幻灯片（在文本框下方）
                                img_top = current_annotation_bottom + \
                                    Inches(0.1)
                                pic = slide.shapes.add_picture(
                                    temp_file_path, panel_start_left, img_top)

                                # 调整图片大小
                                max_img_width = panel_width - Inches(0.2)
                                max_img_height = Inches(1.5)

                                if pic.width > 0 and pic.height > 0:
                                    if pic.width > max_img_width or pic.height > max_img_height:
                                        ratio = min(
                                            max_img_width / pic.width, max_img_height / pic.height)
                                        pic.width = int(pic.width * ratio)
                                        pic.height = int(pic.height * ratio)

                                current_annotation_bottom = img_top + pic.height

                                # 清理临时文件
                                try:
                                    os.unlink(temp_file_path)
                                except OSError:
                                    pass

                                self.logger.info(
                                    f"Added attachment image for annotation {i}")
                    except Exception as e:
                        self.logger.warning(
                            f"Failed to add attachment image for annotation {i}: {e}")

                # 更新下一个注释的位置
                current_top = current_annotation_bottom + Inches(0.2)  # 间距
                annotation_count += 1

            except Exception as e:
                self.logger.warning(
                    f"Failed to add annotation {i} to right panel: {e}")
                continue

        return annotation_count

    async def _apply_drawing_to_image(self, slide, annotation, largest_image):
        """
        将绘画数据应用到图片上，生成带有手绘内容的合成图像

        Args:
            slide: PPT幻灯片对象
            annotation: 包含绘画数据的注释对象
            largest_image: 幻灯片上的最大图像对象
        """
        if not largest_image:
            self.logger.warning("No image found to apply drawing data")
            return

        # 获取绘画数据
        drawing_data = None
        if hasattr(annotation, 'drawing_data'):
            drawing_data = annotation.drawing_data
        elif isinstance(annotation, dict):
            drawing_data = annotation.get('drawing_data')

        if not drawing_data:
            return

        self.logger.info(f"Processing drawing data for annotation")

        try:

            # 解析fabric.js的JSON数据
            fabric_data = json.loads(drawing_data)
            self.logger.info(
                f"Parsed fabric data keys: {list(fabric_data.keys())}")

            if 'objects' in fabric_data:
                self.logger.info(
                    f"Found {len(fabric_data['objects'])} fabric objects")
                for i, obj in enumerate(fabric_data['objects']):
                    obj_type = obj.get('type')
                    self.logger.info(
                        f"Object {i}: type={obj_type}, keys={list(obj.keys())}")

            # 获取原图像的位置和尺寸信息
            img_left_emu = largest_image.left
            img_top_emu = largest_image.top
            img_width_emu = largest_image.width
            img_height_emu = largest_image.height

            # 从图像形状中获取原始图像字节
            original_image_bytes = await self._get_image_bytes_from_shape(largest_image)
            if not original_image_bytes:
                self.logger.warning(
                    "Could not extract original image bytes for drawing overlay - skipping drawing application")
                return

            # 创建PIL图像
            try:
                original_image = Image.open(io.BytesIO(original_image_bytes))
                self.logger.info(
                    f"Successfully loaded original image: size={original_image.size}, mode={original_image.mode}")
            except Exception as e:
                self.logger.error(
                    f"Failed to open original image from bytes: {e}")
                return

            # 创建绘画图层
            drawing_layer = Image.new(
                'RGBA', original_image.size, (0, 0, 0, 0))
            draw = ImageDraw.Draw(drawing_layer)

            # 处理fabric.js对象
            if 'objects' in fabric_data:
                for obj in fabric_data['objects']:
                    await self._render_fabric_object(draw, obj, original_image.size)

            # 合成图像（原图 + 绘画图层）
            if original_image.mode != 'RGBA':
                original_image = original_image.convert('RGBA')

            composite_image = Image.alpha_composite(
                original_image, drawing_layer)

            # 转换回RGB如果需要
            if composite_image.mode == 'RGBA':
                composite_image = composite_image.convert('RGB')

            # 保存合成图像到临时文件
            with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as temp_file:
                composite_image.save(temp_file, format='JPEG', quality=95)
                temp_file_path = temp_file.name

            try:
                # 删除原来的图像
                slide.shapes.element.remove(largest_image.element)

                # 添加合成后的图像到同样的位置
                new_pic = slide.shapes.add_picture(
                    temp_file_path, img_left_emu, img_top_emu)
                new_pic.width = img_width_emu
                new_pic.height = img_height_emu

                self.logger.info("Successfully applied drawing data to image")

            finally:
                # 清理临时文件
                if os.path.exists(temp_file_path):
                    os.unlink(temp_file_path)

        except Exception as e:
            self.logger.error(
                f"Failed to apply drawing data to image: {e}", exc_info=True)

    async def _get_image_bytes_from_shape(self, image_shape):
        """
        从图像形状中提取原始图像字节数据

        Args:
            image_shape: 图像形状对象

        Returns:
            bytes: 图像字节数据，失败时返回None
        """
        try:
            # 尝试从图像形状获取图像数据
            if hasattr(image_shape, 'image') and hasattr(image_shape.image, 'blob'):
                return image_shape.image.blob

            self.logger.warning(
                "Could not extract image bytes from shape - shape may not contain accessible image data")
            return None

        except Exception as e:
            self.logger.error(f"Error extracting image bytes from shape: {e}")
            return None

    async def _render_fabric_object(self, draw, fabric_obj, image_size):
        """
        渲染单个fabric.js对象到PIL绘制上下文

        Args:
            draw: PIL ImageDraw对象
            fabric_obj: fabric.js对象数据
            image_size: 图像尺寸 (width, height)
        """
        try:
            obj_type = fabric_obj.get('type')
            obj_type_lower = obj_type.lower() if obj_type else ''

            if obj_type_lower in ['path']:
                # 处理手绘路径
                await self._render_fabric_path(draw, fabric_obj, image_size)
            elif obj_type_lower in ['circle']:
                # 处理圆形
                await self._render_fabric_circle(draw, fabric_obj, image_size)
            elif obj_type_lower in ['rect', 'rectangle']:
                # 处理矩形
                await self._render_fabric_rect(draw, fabric_obj, image_size)
            elif obj_type_lower in ['line']:
                # 处理直线
                await self._render_fabric_line(draw, fabric_obj, image_size)
            elif obj_type_lower in ['text', 'i-text']:
                # 处理文本
                await self._render_fabric_text(draw, fabric_obj, image_size)
            else:
                self.logger.warning(
                    f"Unsupported fabric.js object type: {obj_type}")

        except Exception as e:
            self.logger.error(f"Error rendering fabric object {obj_type}: {e}")

    async def _render_fabric_path(self, draw, path_obj, image_size):
        """渲染fabric.js路径对象（手绘线条）"""
        try:
            # 获取路径数据
            path_data = path_obj.get('path')
            if not path_data:
                self.logger.warning("No path data found in fabric path object")
                return

            # 获取颜色和线宽
            stroke = path_obj.get('stroke', '#000000')
            stroke_width = max(1, int(path_obj.get('strokeWidth', 2)))

            # 转换颜色
            stroke_color = self._parse_fabric_color(stroke)
            if not stroke_color:
                stroke_color = (0, 0, 0)  # 默认黑色

            self.logger.info(
                f"Rendering path with {len(path_data) if isinstance(path_data, list) else 'unknown'} commands, stroke: {stroke}, width: {stroke_width}")

            # 解析路径并绘制
            points = []

            # 处理路径数据 - fabric.js格式: [["M", x, y], ["Q", x1, y1, x2, y2], ...]
            if isinstance(path_data, list):
                for command in path_data:
                    if not command or not isinstance(command, list) or len(command) < 1:
                        continue

                    cmd = command[0]

                    if cmd == 'M' and len(command) >= 3:  # Move to (绝对坐标)
                        x, y = float(command[1]), float(command[2])
                        points = [(x, y)]

                    elif cmd == 'L' and len(command) >= 3:  # Line to (绝对坐标)
                        x, y = float(command[1]), float(command[2])
                        points.append((x, y))

                    # Quadratic Bezier curve
                    elif cmd == 'Q' and len(command) >= 5:
                        # fabric.js的Q命令格式: ["Q", x1, y1, x2, y2]
                        # x1,y1是控制点，x2,y2是终点
                        # 为简化处理，我们直接连接到终点
                        x, y = float(command[3]), float(command[4])
                        points.append((x, y))

                    elif cmd == 'C' and len(command) >= 7:  # Cubic Bezier curve
                        # fabric.js的C命令格式: ["C", x1, y1, x2, y2, x3, y3]
                        # x1,y1和x2,y2是控制点，x3,y3是终点
                        # 为简化处理，我们直接连接到终点
                        x, y = float(command[5]), float(command[6])
                        points.append((x, y))

                    elif cmd == 'Z':  # Close path
                        # 闭合路径，连接到起始点
                        if points:
                            points.append(points[0])

            # 绘制连接的线段
            if len(points) > 1:
                self.logger.info(
                    f"Drawing path with {len(points)} points, {len(points)-1} line segments")

                # 对于手绘路径，使用连续线段绘制以保证平滑性
                if len(points) > 2:
                    # 绘制连续的线段路径
                    for i in range(len(points) - 1):
                        start_point = points[i]
                        end_point = points[i + 1]

                        # 使用整数坐标避免PIL的浮点数问题
                        start_int = (int(round(start_point[0])), int(
                            round(start_point[1])))
                        end_int = (int(round(end_point[0])), int(
                            round(end_point[1])))

                        draw.line([start_int, end_int],
                                  fill=stroke_color, width=stroke_width)
                else:
                    # 只有两个点的简单直线
                    start_int = (int(round(points[0][0])), int(
                        round(points[0][1])))
                    end_int = (int(round(points[1][0])), int(
                        round(points[1][1])))
                    draw.line([start_int, end_int],
                              fill=stroke_color, width=stroke_width)

            else:
                self.logger.warning(
                    f"Not enough points to draw path: {len(points)} points")

        except Exception as e:
            self.logger.error(
                f"Error rendering fabric path: {e}", exc_info=True)

    async def _render_fabric_circle(self, draw, circle_obj, image_size):
        """渲染fabric.js圆形对象"""
        try:
            left = float(circle_obj.get('left', 0))
            top = float(circle_obj.get('top', 0))
            radius = float(circle_obj.get('radius', 10))

            # 获取样式
            fill = circle_obj.get('fill', 'transparent')
            stroke = circle_obj.get('stroke')
            stroke_width = int(circle_obj.get('strokeWidth', 1))

            # 计算边界框
            bbox = [left - radius, top - radius, left + radius, top + radius]

            # 绘制填充
            if fill and fill != 'transparent':
                fill_color = self._parse_fabric_color(fill)
                if fill_color:
                    draw.ellipse(bbox, fill=fill_color)

            # 绘制边框
            if stroke:
                stroke_color = self._parse_fabric_color(stroke)
                if stroke_color:
                    draw.ellipse(bbox, outline=stroke_color,
                                 width=stroke_width)

        except Exception as e:
            self.logger.error(f"Error rendering fabric circle: {e}")

    async def _render_fabric_rect(self, draw, rect_obj, image_size):
        """渲染fabric.js矩形对象"""
        try:
            left = float(rect_obj.get('left', 0))
            top = float(rect_obj.get('top', 0))
            width = float(rect_obj.get('width', 0))
            height = float(rect_obj.get('height', 0))

            # 获取样式
            fill = rect_obj.get('fill', 'transparent')
            stroke = rect_obj.get('stroke')
            stroke_width = int(rect_obj.get('strokeWidth', 1))

            # 计算边界框
            bbox = [left, top, left + width, top + height]

            # 绘制填充
            if fill and fill != 'transparent':
                fill_color = self._parse_fabric_color(fill)
                if fill_color:
                    draw.rectangle(bbox, fill=fill_color)

            # 绘制边框
            if stroke:
                stroke_color = self._parse_fabric_color(stroke)
                if stroke_color:
                    draw.rectangle(bbox, outline=stroke_color,
                                   width=stroke_width)

        except Exception as e:
            self.logger.error(f"Error rendering fabric rectangle: {e}")

    async def _render_fabric_line(self, draw, line_obj, image_size):
        """渲染fabric.js直线对象"""
        try:
            x1 = float(line_obj.get('x1', 0))
            y1 = float(line_obj.get('y1', 0))
            x2 = float(line_obj.get('x2', 0))
            y2 = float(line_obj.get('y2', 0))
            left = float(line_obj.get('left', 0))
            top = float(line_obj.get('top', 0))

            # 调整坐标
            start_point = (x1 + left, y1 + top)
            end_point = (x2 + left, y2 + top)

            # 获取样式
            stroke = line_obj.get('stroke', '#000000')
            stroke_width = int(line_obj.get('strokeWidth', 1))

            stroke_color = self._parse_fabric_color(stroke)
            if stroke_color:
                draw.line([start_point, end_point],
                          fill=stroke_color, width=stroke_width)

        except Exception as e:
            self.logger.error(f"Error rendering fabric line: {e}")

    async def _render_fabric_text(self, draw, text_obj, image_size):
        """渲染fabric.js文本对象"""
        try:
            text = text_obj.get('text', '')
            if not text:
                return

            left = float(text_obj.get('left', 0))
            top = float(text_obj.get('top', 0))

            # 获取样式
            fill = text_obj.get('fill', '#000000')
            font_size = int(text_obj.get('fontSize', 20))

            fill_color = self._parse_fabric_color(fill)
            if fill_color:
                # 使用默认字体
                try:
                    font = ImageFont.load_default()
                    draw.text((left, top), text, fill=fill_color, font=font)
                except:
                    # 如果字体加载失败，使用默认绘制
                    draw.text((left, top), text, fill=fill_color)

        except Exception as e:
            self.logger.error(f"Error rendering fabric text: {e}")

    def _parse_fabric_color(self, color_value):
        """解析fabric.js颜色值为RGB元组"""
        try:
            if not color_value:
                return None

            if color_value.startswith('#'):
                # 十六进制颜色
                color_value = color_value[1:]
                if len(color_value) == 6:
                    return tuple(int(color_value[i:i+2], 16) for i in (0, 2, 4))
            elif color_value.startswith('rgb'):
                # RGB颜色
                values = re.findall(r'\d+', color_value)
                if len(values) >= 3:
                    return tuple(int(v) for v in values[:3])

            return (0, 0, 0)  # 默认黑色

        except Exception:
            return (0, 0, 0)  # 默认黑色
